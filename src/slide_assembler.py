# src/slide_assembler.py (updated)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os


def assemble_slides(slides_data, output_path):
    prs = Presentation()

    # Use blank layout to avoid placeholder issues
    blank_layout = prs.slide_layouts[6]  # Reliable blank index

    # Light theme background (Alice blue)
    background = prs.slide_master.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)  # Soft blue

    for i, slide_data in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)

        # Title: Top, centered, larger font
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data['title']
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)  # Navy
        p.alignment = PP_ALIGN.CENTER

        # Bullets: Left side, mid-vertical
        bullets_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5), Inches(4))
        tf = bullets_box.text_frame
        tf.word_wrap = True
        for bullet in slide_data['bullets']:
            p = tf.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(20)
            p.level = 0  # Adjust level for sub-bullets if needed

        # Visual: Right side
        visual_path = slide_data.get('visual_path', None)
        if visual_path and os.path.exists(visual_path):
            slide.shapes.add_picture(visual_path, Inches(5.5), Inches(1.5), width=Inches(4), height=Inches(3.5))
        else:
            # Fallback placeholder text if image missing
            fallback = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(3.5))
            tf = fallback.text_frame
            p = tf.paragraphs[0]
            p.text = "Visual Placeholder\n(No image found)"
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER

        # Footer: Slide number + tagline
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = f"Slide {i} | KNN Explainer Prototype"
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER

    prs.save(output_path)